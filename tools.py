import os
import subprocess
import difflib
import webbrowser
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

def get_file_diff(path: str, new_content: str) -> str:
    """
    Generates a unified diff between the current file content and the new content.
    """
    try:
        if not os.path.exists(path):
            return f"--- (New File: {path}) ---\n{new_content}"
        
        with open(path, 'r', encoding='utf-8') as f:
            old_content = f.read()
            
        diff = difflib.unified_diff(
            old_content.splitlines(keepends=True),
            new_content.splitlines(keepends=True),
            fromfile=f"a/{path}",
            tofile=f"b/{path}"
        )
        return "".join(diff)
    except Exception as e:
        return f"Error generating diff: {str(e)}"

def run_command(command: str, cwd: str = None, wait: bool = None) -> str:
    """
    Executes a shell command. 
    - wait: If True, blocks until completion. If False, runs in background.
    By default, it auto-detects servers/GUIs to run in background.
    """
    # Detect common server/blocking keywords
    blocking_keywords = ["runserver", "npm start", "npm run dev", "vite", "serve", "watch", "python app.py", "python main.py"]
    is_server = any(kw in command.lower() for kw in blocking_keywords)
    
    # If agent explicitly asks not to wait, or if it's a server, use background
    should_wait = wait if wait is not None else (not is_server)
    
    exec_cwd = cwd if (cwd and os.path.exists(cwd)) else os.getcwd()

    try:
        if not should_wait:
            # Run in background using Popen
            process = subprocess.Popen(
                command,
                shell=True,
                cwd=exec_cwd,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                creationflags=subprocess.CREATE_NEW_PROCESS_GROUP if os.name == 'nt' else 0
            )
            return f"Started in background (PID: {process.pid}) at {exec_cwd}. The process is running independently."

        # Run command synchronously and capture output
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            cwd=exec_cwd,
            timeout=60 # Safety timeout for synchronous commands
        )
        output = result.stdout
        if result.stderr:
            output += "\n\n--- stderr ---\n\n" + result.stderr
        
        if result.returncode != 0:
            return f"Command failed (Code {result.returncode}) at {exec_cwd}\nOutput:\n{output}"
            
        return f"Command success in {exec_cwd}.\nOutput:\n{output}"
    except subprocess.TimeoutExpired:
        return f"Command timed out after 60s in {exec_cwd}. It might be a persistent process; try running with wait=False."
    except Exception as e:
        return f"Tool Error: {str(e)}"

def git_init(cwd: str = None) -> str:
    """Initializes a new git repository in the project root."""
    return run_command("git init", cwd=cwd)

def git_commit(message: str, cwd: str = None) -> str:
    """Stages all changes and creates a commit."""
    run_command("git add .", cwd=cwd)
    return run_command(f'git commit -m "{message}"', cwd=cwd)

def git_push(remote: str = "origin", branch: str = "main", cwd: str = None) -> str:
    """Pushes local commits to a remote repository."""
    return run_command(f"git push {remote} {branch}", cwd=cwd)

def project_auto_analyze(cwd: str = None) -> str:
    """
    Automatically scans core project files (README, requirements, app.py) 
    to determine the tech stack and purpose.
    """
    base = cwd if (cwd and os.path.exists(cwd)) else os.getcwd()
    analysis = [f"Project Analysis for {base}:"]
    
    critical_files = ["README.md", "README", "requirements.txt", "package.json", "app.py", "main.py", "manage.py"]
    for f in critical_files:
        p = os.path.join(base, f)
        if os.path.exists(p):
            analysis.append(f"\n--- Found {f} ---")
            try:
                with open(p, 'r', encoding='utf-8') as file:
                    content = file.read(1000) # Read first 1000 chars
                    analysis.append(content + ("..." if len(content) == 1000 else ""))
            except Exception as e:
                analysis.append(f"[File unreadable or locked: {str(e)}]")
            
    if len(analysis) == 1:
        return f"No core files found in {base}. Directory seems empty or unstructured."
        
    return "\n".join(analysis)

def write_file(path: str, content: str, cwd: str = None) -> str:
    """
    Writes content to a file. Path should be relative to the project root.
    """
    try:
        target_base = cwd if (cwd and os.path.exists(cwd)) else os.getcwd()
        safe_path = path.lstrip('/\\')
        full_path = os.path.join(target_base, safe_path)
        
        os.makedirs(os.path.dirname(os.path.abspath(full_path)), exist_ok=True)
        with open(full_path, 'w', encoding='utf-8') as f:
            f.write(content)
        return f"Successfully wrote to {full_path}"
    except Exception as e:
        return f"Failed to write file {path}: {str(e)}"

def read_file(path: str, cwd: str = None) -> str:
    """
    Reads the content of a file. Path should be relative to the project root.
    """
    try:
        target_base = cwd if (cwd and os.path.exists(cwd)) else os.getcwd()
        safe_path = path.lstrip('/\\')
        full_path = os.path.join(target_base, safe_path)
        
        if not os.path.exists(full_path):
            return f"File not found: {full_path}"
        with open(full_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"Failed to read file {path}: {str(e)}"

def list_directory(path: str = ".", cwd: str = None) -> str:
    """
    Lists the contents of a directory. Path should be relative to the project root.
    """
    try:
        # Resolve target path relative to provided cwd
        target_base = cwd if (cwd and os.path.exists(cwd)) else os.getcwd()
        safe_path = path.lstrip('/\\')
        full_path = os.path.join(target_base, safe_path)
        
        if not os.path.exists(full_path):
            return f"Path not found: {full_path}"
        
        items = os.listdir(full_path)
        output_lines = [f"Contents of {os.path.abspath(full_path)}:"]
        for item in items:
            item_path = os.path.join(full_path, item)
            item_type = "(dir)" if os.path.isdir(item_path) else "(file)"
            output_lines.append(f"  {item} {item_type}")
            
        return "\n".join(output_lines)
    except Exception as e:
        return f"Failed to list directory {path}: {str(e)}"

def launch_browser(url: str) -> str:
    """
    Opens the specified URL in the default system web browser.
    Useful for launching the web applications you build (e.g. http://localhost:3000)
    so the user can see them immediately.
    """
    try:
        # webbrowser.open returns True if successful
        success = webbrowser.open(url)
        if success:
            return f"Successfully launched {url} in the default browser."
        else:
            return f"Failed to launch {url}. The browser might not be accessible."
    except Exception as e:
        return f"Error launching browser for {url}: {str(e)}"


def _hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def create_presentation(title: str, slides: list, output_path: str) -> str:
    """
    Creates a professional PowerPoint presentation with advanced styling and flow diagram support.
    - title: Main title of the presentation.
    - slides: A list of dicts:
        {
            "title": "Slide Title",
            "content": ["points"],
            "background_color": "#RRGGBB",
            "text_color": "#RRGGBB",
            "flow_elements": [{"type": "box", "text": "Step 1", "x": 1, "y": 1}, {"type": "arrow", "from_x": 2, "from_y": 1.5}]
        }
    - output_path: .pptx save path.
    """
    try:
        prs = Presentation()
        
        # Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Autonomous WebDev Agent Premium Presentation"

        for s_data in slides:
            # Use blank layout for better manual control of flow diagrams if they exist
            has_flow = "flow_elements" in s_data
            layout_idx = 6 if has_flow else 1 # 6 is Blank, 1 is Title/Content
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            
            # Apply Colors
            bg_hex = s_data.get("background_color", "#FFFFFF")
            txt_hex = s_data.get("text_color", "#000000")
            
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(bg_hex)

            # Add Title if layout is blank
            if has_flow:
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
                tf = title_shape.text_frame
                p = tf.paragraphs[0]
                p.text = s_data.get("title", "")
                p.font.bold = True
                p.font.size = Pt(32)
                p.font.color.rgb = _hex_to_rgb(txt_hex)
            else:
                title_place = slide.shapes.title
                title_place.text = s_data.get("title", "Untitled Slide")
                title_place.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(txt_hex)
                
                body_place = slide.placeholders[1]
                tf = body_place.text_frame
                tf.clear()
                content = s_data.get("content", [])
                if isinstance(content, list):
                    for i, point in enumerate(content):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = point
                        p.font.color.rgb = _hex_to_rgb(txt_hex)
                else:
                    tf.paragraphs[0].text = str(content)
                    tf.paragraphs[0].font.color.rgb = _hex_to_rgb(txt_hex)

            # Add Flow Elements
            if has_flow:
                for elem in s_data.get("flow_elements", []):
                    etype = elem.get("type")
                    x, y = elem.get("x", 0), elem.get("y", 0)
                    w, h = elem.get("w", 2), elem.get("h", 1)
                    
                    if etype == "box":
                        # MSO_SHAPE_TYPE 1 = RECTANGLE
                        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                        shape.text = elem.get("text", "")
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = _hex_to_rgb(txt_hex)
                        shape.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(bg_hex)
                    elif etype == "arrow":
                        fx, fy = elem.get("from_x", 0), elem.get("from_y", 0)
                        tx, ty = elem.get("to_x", 2), elem.get("to_y", 0)
                        
                        import math
                        w = math.sqrt((tx - fx)**2 + (ty - fy)**2)
                        h = 0.4
                        
                        arrow = slide.shapes.add_shape(13, Inches(fx), Inches(fy), Inches(w), Inches(h))
                        arrow.rotation = int(math.degrees(math.atan2(ty - fy, tx - fx)))
                        arrow.fill.solid()
                        arrow.fill.fore_color.rgb = _hex_to_rgb(txt_hex)

            # Add Image if provided
            image_path = s_data.get("image_path")
            if image_path and os.path.exists(image_path):
                # Put image on the right or center depending on flow
                ix = Inches(5.5) if not has_flow else Inches(1)
                iy = Inches(1.5) if not has_flow else Inches(3)
                iw = Inches(4)
                slide.shapes.add_picture(image_path, ix, iy, width=iw)

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)
        return f"Successfully created enhanced presentation at {output_path}"
    except Exception as e:
        return f"Failed to create presentation: {str(e)}"
